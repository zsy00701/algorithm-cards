#!/usr/bin/env python3
"""把 Hot100 方法卡的结构化源数据重建为【原生可编辑】PPTX。

每页所有文字、代码、形状都是原生 PPT 元素，可在 PowerPoint/Keynote 里点选编辑。
源数据直接复用 hot100_method_cards/render_method_cards.py 里的 CARDS。
"""
import importlib.util
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTDIR = os.path.join(ROOT, "ppt_output")
os.makedirs(OUTDIR, exist_ok=True)

# --- 复用源数据 ---
spec = importlib.util.spec_from_file_location(
    "render_method_cards",
    os.path.join(ROOT, "hot100_method_cards", "render_method_cards.py"),
)
rmc = importlib.util.module_from_spec(spec)
spec.loader.exec_module(rmc)
CARDS = rmc.CARDS

# --- 配色 ---
NAVY = RGBColor(0x0B, 0x2E, 0x6D)
BLUE = RGBColor(0x17, 0x69, 0xD2)
CYAN = RGBColor(0x2D, 0x95, 0xE8)
PALE = RGBColor(0xE8, 0xF3, 0xFF)
PALE2 = RGBColor(0xF4, 0xF9, 0xFF)
GOLD = RGBColor(0xE9, 0xA5, 0x1D)
GREEN = RGBColor(0x15, 0x84, 0x5B)
TEXT = RGBColor(0x17, 0x24, 0x3A)
MUTED = RGBColor(0x60, 0x70, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODEBG = RGBColor(0x0E, 0x1B, 0x33)
CODEFG = RGBColor(0xE6, 0xEF, 0xFF)

ZH = "PingFang SC"
MONO = "Menlo"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def _set_ea(run, font_name):
    """让中文也用指定字体（设置 East Asian 字体）。"""
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn("a:ea"))
    if ea is None:
        ea = rpr.makeelement(qn("a:ea"), {})
        rpr.append(ea)
    ea.set("typeface", font_name)


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    return shp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tf


def para(tf, text, size, color, *, bold=False, font=ZH, align=PP_ALIGN.LEFT,
         first=False, space_before=0, space_after=2, line=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = font
    _set_ea(r, font)
    return p


def section_head(slide, x, y, w, label, color=BLUE):
    """小节标题：左侧色条 + 粗体标题。"""
    rect(slide, x, y + 0.02, 0.07, 0.26, fill=color)
    tf = textbox(slide, x + 0.16, y, w - 0.16, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, label, 13, NAVY, bold=True, first=True)


def code_block(slide, x, y, w, h, lang, code, accent):
    rect(slide, x, y, w, h, fill=CODEBG, rounded=True)
    # 语言标签条
    rect(slide, x, y, w, 0.30, fill=accent, rounded=True)
    rect(slide, x, y + 0.14, w, 0.16, fill=accent)  # 盖住下圆角，做成上圆下方
    tf_l = textbox(slide, x + 0.12, y, w - 0.2, 0.30, anchor=MSO_ANCHOR.MIDDLE)
    para(tf_l, lang, 11, WHITE, bold=True, font=MONO, first=True)
    # 代码体
    tf = textbox(slide, x + 0.05, y + 0.34, w - 0.1, h - 0.4)
    tf.word_wrap = False
    lines = code.split("\n")
    for i, ln in enumerate(lines):
        para(tf, ln if ln else " ", 8.5, CODEFG, font=MONO,
             first=(i == 0), space_after=0, line=1.02)
    return tf


# ---------------- 封面 ----------------
def cover():
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, fill=PALE2)
    rect(s, 0, 0, 4.7, 7.5, fill=NAVY)
    rect(s, 0.55, 0.7, 1.0, 0.09, fill=CYAN)
    tf = textbox(s, 0.55, 0.95, 3.7, 0.5)
    para(tf, "ALGORITHM MAP", 16, RGBColor(0x8B, 0xC2, 0xFF), bold=True, font="Arial", first=True)
    tf = textbox(s, 0.55, 1.7, 3.8, 2.2)
    para(tf, "力扣 Hot 100\n算法方法图解", 40, WHITE, bold=True, first=True, line=1.1)
    tf = textbox(s, 0.55, 4.0, 3.7, 1.0)
    para(tf, "18 个高频方法专题", 17, RGBColor(0xD8, 0xE9, 0xFF), bold=True, first=True)
    para(tf, "C++ / Python 双语代码模板 · 原生可编辑版", 14, RGBColor(0xB8, 0xD4, 0xF5))
    tf = textbox(s, 0.55, 6.5, 3.7, 0.6)
    para(tf, "", 13, RGBColor(0xB8, 0xD4, 0xF5), first=True)
    # 右侧目录预览（前 9 个标题）
    tf = textbox(s, 5.2, 0.9, 7.6, 0.6)
    para(tf, "目录速览", 22, NAVY, bold=True, first=True)
    col_w = 3.75
    for i, c in enumerate(CARDS):
        col = i // 9
        row = i % 9
        tf = textbox(s, 5.2 + col * col_w, 1.7 + row * 0.58, col_w - 0.1, 0.55,
                     anchor=MSO_ANCHOR.MIDDLE)
        p = para(tf, f"{i+1:02d}  ", 13, BLUE, bold=True, font=MONO, first=True)
        r = p.add_run()
        r.text = c["title"]
        r.font.size = Pt(13)
        r.font.color.rgb = TEXT
        r.font.name = ZH
        _set_ea(r, ZH)


# ---------------- 卡片页 ----------------
def card_slide(card, index):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, 13.333, 7.5, fill=WHITE)

    # 顶部标题栏
    rect(s, 0, 0, 13.333, 0.92, fill=NAVY)
    rect(s, 0, 0.92, 13.333, 0.06, fill=CYAN)
    tf = textbox(s, 0.45, 0, 10.5, 0.92, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, card["title"], 24, WHITE, bold=True, first=True)
    badge = rect(s, 11.95, 0.26, 1.0, 0.42, fill=BLUE, rounded=True)
    btf = badge.text_frame
    btf.word_wrap = False
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    br = bp.add_run()
    br.text = f"{index:02d} / 18"
    br.font.size = Pt(13)
    br.font.bold = True
    br.font.color.rgb = WHITE
    br.font.name = MONO

    LX, LW = 0.45, 6.1
    RX, RW = 6.9, 5.95
    y = 1.18

    # 定义
    section_head(s, LX, y, LW, "定义")
    tf = textbox(s, LX, y + 0.34, LW, 0.75)
    para(tf, card["definition"], 13, TEXT, first=True, line=1.12)

    # 识别信号
    y2 = 2.28
    section_head(s, LX, y2, LW, "识别信号")
    tf = textbox(s, LX, y2 + 0.34, LW, 1.0)
    for i, sig in enumerate(card["signals"]):
        p = para(tf, "●  ", 11, CYAN, first=(i == 0), space_after=3)
        r = p.add_run()
        r.text = sig
        r.font.size = Pt(12.5)
        r.font.color.rgb = TEXT
        r.font.name = ZH
        _set_ea(r, ZH)

    # 解题步骤
    y3 = 3.62
    section_head(s, LX, y3, LW, "解题步骤")
    tf = textbox(s, LX, y3 + 0.34, LW, 1.7)
    for i, st in enumerate(card["steps"]):
        p = para(tf, f"{i+1}. ", 12.5, BLUE, bold=True, font=MONO,
                 first=(i == 0), space_after=4)
        r = p.add_run()
        r.text = st
        r.font.size = Pt(12.5)
        r.font.color.rgb = TEXT
        r.font.name = ZH
        _set_ea(r, ZH)

    # 例题（左下高亮条）
    y4 = 5.95
    rect(s, LX, y4, LW, 0.62, fill=PALE, line=RGBColor(0xB7, 0xD5, 0xFA), rounded=True)
    tf = textbox(s, LX + 0.18, y4, LW - 0.3, 0.62, anchor=MSO_ANCHOR.MIDDLE)
    p = para(tf, "例  ", 11, GOLD, bold=True, first=True)
    r = p.add_run()
    r.text = card["example"]
    r.font.size = Pt(11.5)
    r.font.color.rgb = NAVY
    r.font.name = ZH
    _set_ea(r, ZH)

    # 右列：C++ / Python 代码
    section_head(s, RX, 1.18, RW, "C++ 模板", color=NAVY)
    code_block(s, RX, 1.52, RW, 2.0, "C++", card["cpp"], NAVY)
    section_head(s, RX, 3.66, RW, "Python 模板", color=BLUE)
    code_block(s, RX, 4.0, RW, 1.95, "Python", card["py"], BLUE)

    # 右下：记忆口诀 + Hot100 题号
    yb = 6.05
    rect(s, RX, yb, RW, 0.62, fill=RGBColor(0xFF, 0xF6, 0xE5),
         line=RGBColor(0xF0, 0xD8, 0xA8), rounded=True)
    tf = textbox(s, RX + 0.18, yb + 0.04, RW - 0.3, 0.55, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, card["memory"], 11.5, RGBColor(0x9A, 0x6A, 0x10), bold=True, first=True)

    tf = textbox(s, RX, yb + 0.66, RW, 0.42)
    para(tf, card["hot"], 9.5, MUTED, first=True, line=1.05)


cover()
for i, c in enumerate(CARDS, 1):
    card_slide(c, i)

out = os.path.join(OUTDIR, "Hot100方法卡_可编辑版.pptx")
prs.save(out)
print(f"[完成] {out}  共 {len(prs.slides.__iter__.__self__._sldIdLst)} 页")
