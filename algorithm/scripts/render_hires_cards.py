#!/usr/bin/env python3
"""把 hot100 方法卡用 render_method_cards.py 的同款布局，按 S 倍高分辨率重渲染。

做法：复用原模块的全部绘制逻辑（含示例图解 draw_visual），但通过一个
ScaledDraw 代理把所有坐标、线宽、半径、字号统一放大 S 倍——
文字/代码/形状都在 S 倍画布上原生绘制 → 清晰；几何测量仍在 1x 空间 → 不串版。
内容来自源数据，保证正确。
"""
import importlib.util
import os

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
S = 2  # 放大倍数

# --- 载入原渲染模块 ---
spec = importlib.util.spec_from_file_location(
    "rmc", os.path.join(ROOT, "hot100_method_cards", "render_method_cards.py"))
M = importlib.util.module_from_spec(spec)
spec.loader.exec_module(M)

# --- 字体注册（记录 path/size/index，供代理放大重建）---
_FR = {}
_KEEP = []
_orig_zh, _orig_mono = M.zh, M.mono


def zh2(size, bold=False):
    f = _orig_zh(size, bold)
    _FR[id(f)] = (M.ZH_BOLD if bold else M.ZH_REG, size, 0)
    _KEEP.append(f)
    return f


def mono2(size, bold=False):
    f = _orig_mono(size, bold)
    _FR[id(f)] = (M.MONO, size, 1 if bold else 0)
    _KEEP.append(f)
    return f


M.zh, M.mono = zh2, mono2


class ScaledDraw:
    """坐标/线宽/半径/字号 ×S 后转发给真实 draw；文字测量保持 1x。"""

    def __init__(self, real, s):
        self.real, self.s = real, s
        self._fontcache = {}

    def _sc(self, o):
        if isinstance(o, bool):
            return o
        if isinstance(o, (int, float)):
            return o * self.s
        if isinstance(o, (list, tuple)):
            return type(o)(self._sc(x) for x in o)
        return o

    def _w(self, width):
        return max(1, int(round(width * self.s)))

    def _font(self, font):
        if font is None:
            return None
        key = id(font)
        path, size, index = _FR[key]
        ck = (path, int(round(size * self.s)), index)
        if ck not in self._fontcache:
            self._fontcache[ck] = ImageFont.truetype(ck[0], ck[1], index=ck[2])
        return self._fontcache[ck]

    def text(self, xy, text, font=None, fill=None, **kw):
        self.real.text(self._sc(xy), text, font=self._font(font), fill=fill, **kw)

    def line(self, xy, fill=None, width=1, **kw):
        self.real.line(self._sc(xy), fill=fill, width=self._w(width), **kw)

    def rectangle(self, xy, fill=None, outline=None, width=1, **kw):
        self.real.rectangle(self._sc(xy), fill=fill, outline=outline,
                            width=self._w(width), **kw)

    def rounded_rectangle(self, xy, radius=0, fill=None, outline=None, width=1, **kw):
        self.real.rounded_rectangle(self._sc(xy), radius=radius * self.s, fill=fill,
                                    outline=outline, width=self._w(width), **kw)

    def ellipse(self, xy, fill=None, outline=None, width=1, **kw):
        self.real.ellipse(self._sc(xy), fill=fill, outline=outline,
                         width=self._w(width), **kw)

    def polygon(self, xy, fill=None, outline=None, **kw):
        self.real.polygon(self._sc(xy), fill=fill, outline=outline, **kw)

    def textbbox(self, xy, text, font=None, **kw):
        return self.real.textbbox(xy, text, font=font, **kw)  # 1x 测量


def render_card_hires(card, out_dir):
    base = Image.open(M.BASE).convert("RGBA")
    w, h = base.size
    img = base.resize((w * S, h * S), Image.Resampling.LANCZOS)
    real = ImageDraw.Draw(img, "RGBA")
    draw = ScaledDraw(real, S)
    zh, mono = M.zh, M.mono
    NAVY, BLUE, TEXT, WHITE = M.NAVY, M.BLUE, M.TEXT, M.WHITE

    # —— 与 render_card 完全一致的绘制序列 ——
    M.center(draw, 728, 72, card["title"], zh(48, True), NAVY)
    draw.line((310, 112, 1145, 112), fill=BLUE, width=3)

    draw.text((118, 158), "定义", font=zh(24, True), fill=NAVY)
    M.text_block(draw, (175, 156), card["definition"], zh(20), TEXT, 505, 5, 3)

    draw.text((810, 158), "识别信号", font=zh(24, True), fill=NAVY)
    for i, signal in enumerate(card["signals"]):
        draw.polygon([(820, 191 + i * 28), (832, 198 + i * 28), (820, 205 + i * 28)], fill=BLUE)
        draw.text((842, 181 + i * 28), signal, font=zh(19, True if i == 0 else False),
                  fill=TEXT if i else BLUE)

    draw.text((70, 292), "示例推演", font=zh(25, True), fill=NAVY)
    draw.text((205, 295), card["example"], font=zh(20), fill=TEXT)
    M.draw_visual(draw, card["visual"])

    draw.text((1115, 294), "算法步骤", font=zh(25, True), fill=NAVY)
    for i, step in enumerate(card["steps"]):
        cy = 356 + i * 80
        M.center(draw, 1090, cy, str(i + 1), zh(21, True), WHITE)
        M.text_block(draw, (1120, cy - 20), step, zh(19), TEXT, 255, 5, 2)

    M.draw_code(draw, 48, "C++", card["cpp"], NAVY)
    M.draw_code(draw, 738, "Python", card["py"], BLUE)

    M.center(draw, 728, 966, card["memory"], zh(22, True), NAVY)
    M.center(draw, 728, 1013, card["hot"], zh(17), TEXT)

    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, f"{card['slug']}.png")
    img.convert("RGB").save(path, quality=95)
    return path


def build_pptx(image_paths, out_path):
    w, h = Image.open(image_paths[0]).size
    ratio = w / h
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(13.333 / ratio)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    for p in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(p, 0, 0, width=SW, height=SH)
    prs.save(out_path)


if __name__ == "__main__":
    out_dir = os.path.join(ROOT, "hot100_method_cards", "hires")
    paths = []
    for i, card in enumerate(M.CARDS, 1):
        paths.append(render_card_hires(card, out_dir))
        print(f"  渲染 {i:02d}/18  {card['slug']}")
    w, h = Image.open(paths[0]).size
    print(f"高清卡尺寸: {w}x{h}")
    ppt = os.path.join(ROOT, "ppt_output", "Hot100方法卡_正确高清版.pptx")
    build_pptx(paths, ppt)
    print(f"[完成] {ppt}")
