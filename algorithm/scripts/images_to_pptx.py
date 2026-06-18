#!/usr/bin/env python3
"""把各组生成的图片分别汇总成独立的 PPTX，每张图占满一页。"""
import glob
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTDIR = os.path.join(ROOT, "ppt_output")
os.makedirs(OUTDIR, exist_ok=True)

EMU_PER_INCH = 914400
MAX_INCH = 13.333  # 长边最大英寸数

# (源目录, 输出文件名)
GROUPS = [
    ("carl_algorithm_cards/gpt-image-2",   "代码随想录算法卡_gpt-image-2.pptx"),
    ("hot100_images/output",               "Hot100分类总览.pptx"),
    ("hot100_method_cards/output",         "Hot100方法卡_output.pptx"),
    ("hot100_method_cards/gpt-image-2",    "Hot100方法卡_gpt-image-2.pptx"),
]

# 跳过这些辅助图（预览图/背景图）
SKIP_KEYWORDS = ("preview", "background")


def make_pptx(src_dir, out_name):
    abs_dir = os.path.join(ROOT, src_dir)
    files = sorted(glob.glob(os.path.join(abs_dir, "*.png")))
    files = [f for f in files
             if not any(k in os.path.basename(f).lower() for k in SKIP_KEYWORDS)]
    if not files:
        print(f"[跳过] {src_dir} 无图片")
        return

    # 用首图比例决定页面尺寸（组内尺寸一致）
    w0, h0 = Image.open(files[0]).size
    ratio = w0 / h0
    if ratio >= 1:
        sw_in, sh_in = MAX_INCH, MAX_INCH / ratio
    else:
        sh_in, sw_in = MAX_INCH, MAX_INCH * ratio

    prs = Presentation()
    prs.slide_width = Emu(int(sw_in * EMU_PER_INCH))
    prs.slide_height = Emu(int(sh_in * EMU_PER_INCH))
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    for f in files:
        slide = prs.slides.add_slide(blank)
        iw, ih = Image.open(f).size
        ir = iw / ih
        sr = SW / SH
        if ir > sr:  # 图更宽 -> 按宽撑满
            pw = SW
            ph = int(SW / ir)
        else:        # 图更高 -> 按高撑满
            ph = SH
            pw = int(SH * ir)
        left = int((SW - pw) / 2)
        top = int((SH - ph) / 2)
        slide.shapes.add_picture(f, left, top, width=pw, height=ph)

    out_path = os.path.join(OUTDIR, out_name)
    prs.save(out_path)
    print(f"[完成] {out_name}: {len(files)} 页  ({sw_in:.2f}x{sh_in:.2f} in)")


for src, name in GROUPS:
    make_pptx(src, name)

print(f"\n全部输出到: {OUTDIR}")
